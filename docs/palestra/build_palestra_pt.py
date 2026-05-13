"""
Gerador da palestra Zerbinatti (5min) — slide deck editorial premium — VERSAO PT.
Saida: /Users/fabricio/dev/cafe/docs/palestra/zerbinatti-palestra-5min-PT.pptx
"""
import os
import json
import urllib.request
from io import BytesIO
from PIL import Image, ImageFilter

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.patches import Circle, Ellipse, FancyArrowPatch, PathPatch
from matplotlib.path import Path as MplPath
import matplotlib.patheffects as path_effects
import numpy as np

# -----------------------------------------------------------------------------
# Constantes de identidade visual
# -----------------------------------------------------------------------------
BLACK = RGBColor(0x0A, 0x0A, 0x0A)
CREAM = RGBColor(0xEC, 0xEC, 0xEC)
GOLD = RGBColor(0xC9, 0xA9, 0x61)
GOLD_LIGHT = RGBColor(0xDD, 0xC8, 0xA0)
GOLD_PALE = RGBColor(0xF0, 0xE0, 0xC0)
GOLD_DARK = RGBColor(0x8A, 0x73, 0x42)

# Fontes (sistema macOS)
F_SERIF = "Cormorant Garamond"
F_SERIF_FB = "Garamond"
F_SANS = "Inter"
F_SANS_FB = "Helvetica Neue"

# Caminhos
ROOT = "/Users/fabricio/dev/cafe"
OUT_DIR = f"{ROOT}/docs/palestra"
IMG_DIR = f"{OUT_DIR}/img"
OUT_PATH = f"{OUT_DIR}/zerbinatti-palestra-5min-PT.pptx"

os.makedirs(IMG_DIR, exist_ok=True)

# -----------------------------------------------------------------------------
# Converter webp -> jpg (python-pptx nao le webp)
# -----------------------------------------------------------------------------
WEBP_FILES = {
    "cherries-yellow": f"{ROOT}/public/images/farm/cherries-yellow.webp",
    "drying-yard-tree": f"{ROOT}/public/images/farm/drying-yard-tree.webp",
    "flowers-sunset": f"{ROOT}/public/images/farm/flowers-sunset.webp",
    "young-plants": f"{ROOT}/public/images/farm/young-plants.webp",
    "seedlings-nursery": f"{ROOT}/public/images/farm/seedlings-nursery.webp",
    "packages-row": f"{ROOT}/public/images/farm/packages-row.webp",
    "peneirar": f"{ROOT}/public/assets/galeria/peneirar.webp",
}

JPGS = {}
for name, src in WEBP_FILES.items():
    dst = f"{IMG_DIR}/{name}.jpg"
    if not os.path.exists(dst):
        im = Image.open(src).convert("RGB")
        im.save(dst, "JPEG", quality=92, optimize=True)
    JPGS[name] = dst
    print(f"  jpg ok: {dst}")

WORDMARK_GOLD = f"{ROOT}/public/assets/zerbinatti-wordmark-gold.png"
WORDMARK_CREAM = f"{ROOT}/public/assets/zerbinatti-wordmark-cream.png"
QR_PATH = f"{IMG_DIR}/qr-zerbinatti.png"

# -----------------------------------------------------------------------------
# Geracao do mapa-mundi estilizado em ouro/preto para slide 4 (PT)
# -----------------------------------------------------------------------------
GEO_URL = "https://raw.githubusercontent.com/nvkelso/natural-earth-vector/master/geojson/ne_110m_admin_0_countries.geojson"
GEO_PATH = f"{IMG_DIR}/ne_110m_countries.geojson"


def build_world_map():
    """Mapa-mundi minimalista em ouro/preto usando contornos reais (Natural Earth 110m).
    Marca: Brasil (origem), Franca/Suica/Suecia (presenca), Espanha (proximo), Oregon (sonho).
    Cada destino recebe uma rota tracejada saindo do Brasil.
    """
    out = f"{IMG_DIR}/world-map-pt.png"

    # Cache do GeoJSON Natural Earth (low-res, ~200KB)
    if not os.path.exists(GEO_PATH):
        print(f"  baixando geojson: {GEO_URL}")
        urllib.request.urlretrieve(GEO_URL, GEO_PATH)
    with open(GEO_PATH) as f:
        world = json.load(f)

    fig, ax = plt.subplots(figsize=(16, 8), facecolor="#0A0A0A")
    ax.set_facecolor("#0A0A0A")
    ax.set_xlim(-180, 180)
    ax.set_ylim(-58, 80)
    ax.set_aspect("equal")
    ax.axis("off")

    GOLD_HEX = "#D4B574"     # borda mais clara/luminosa
    GOLD_FILL = "#4A3A1F"    # fill mais quente

    # (A) Atmosfera de fundo (oceano) — gradiente radial sutil
    for r, alpha in [(220, 0.025), (160, 0.04), (100, 0.05)]:
        ax.add_patch(Ellipse((0, 10), width=r * 2, height=r * 1.1,
                              facecolor="#C9A961", edgecolor="none",
                              alpha=alpha, zorder=0))

    # (C) Equador + tropicos (linhas pontilhadas discretas)
    for y_lat, a_line in [(0, 0.25), (23.5, 0.12), (-23.5, 0.12)]:
        ax.axhline(y=y_lat, color="#C9A961", linewidth=0.4,
                   linestyle=(0, (1, 3)), alpha=a_line, zorder=1)

    # (B) Países com profundidade — borda mais grossa e cores mais quentes
    for feat in world["features"]:
        if feat.get("properties", {}).get("NAME") == "Antarctica":
            continue
        geom = feat.get("geometry")
        if geom is None:
            continue
        gtype = geom["type"]
        if gtype == "MultiPolygon":
            polys = geom["coordinates"]
        elif gtype == "Polygon":
            polys = [geom["coordinates"]]
        else:
            continue
        for poly in polys:
            ring = poly[0]
            patch = mpatches.Polygon(
                ring, closed=True,
                facecolor=GOLD_FILL, edgecolor=GOLD_HEX,
                linewidth=0.9, alpha=0.95, joinstyle="round",
                zorder=2,
            )
            ax.add_patch(patch)

    points = [
        (-47, -15, "BRASIL · ORIGEM",  "origin"),
        (2,   47, "FRANÇA",            "present"),
        (8,   46, "SUÍÇA",             "present"),
        (16,  60, "SUÉCIA",            "present"),
        (-4,  40, "ESPANHA · PRÓXIMO", "next"),
        (-122, 45, "OREGON · SONHO",   "dream"),
    ]

    # (F) Rotas com glow — desenha 2x: glow grosso + linha tracejada
    # Cada destino tem altura de arco diferente -> efeito arco-iris empilhado
    ARC_HEIGHT = {
        "ESPANHA · PRÓXIMO": 0.10,
        "FRANÇA":            0.32,
        "SUÍÇA":             0.56,
        "SUÉCIA":            0.82,
        "OREGON · SONHO":    0.22,
    }
    # cor distinta por destino — arco-iris empilhado
    ROUTE_STYLE = {
        "ESPANHA · PRÓXIMO": ("#FF4FA3", 1.8, (0, (6, 4)),   1.00),  # rosa
        "FRANÇA":            ("#FFD93C", 1.5, (0, (5, 4)),   1.00),  # amarelo
        "SUÍÇA":             ("#3CD982", 1.5, (0, (4, 4)),   1.00),  # verde
        "SUÉCIA":            ("#A878FF", 1.5, (0, (3, 4)),   1.00),  # roxo claro
    }
    # rota especial Oregon: tracejado intercalado azul/vermelho (cores EUA)
    OREGON_LABEL = "OREGON · SONHO"
    origin = (-47, -15)
    for lon, lat, label, role in points:
        if role == "origin":
            continue
        mid_x = (origin[0] + lon) / 2
        h = ARC_HEIGHT.get(label, 0.18)
        mid_y = (origin[1] + lat) / 2 + (abs(lon - origin[0]) * h)
        verts = [origin, (mid_x, mid_y), (lon, lat)]
        codes = [MplPath.MOVETO, MplPath.CURVE3, MplPath.CURVE3]
        path = MplPath(verts, codes)

        if label == OREGON_LABEL:
            # intercalar dois tracejados deslocados — vermelho e azul EUA
            lw = 1.6
            # base: azul (offset 0)
            ax.add_patch(PathPatch(
                path, fill=False, edgecolor="#1E5BD8", linewidth=lw * 1.4,
                linestyle=(0, (6, 8)), alpha=1.0, zorder=5,
            ))
            # sobreposto: vermelho (offset 7 = comeca onde o azul termina)
            ax.add_patch(PathPatch(
                path, fill=False, edgecolor="#E63946", linewidth=lw * 1.4,
                linestyle=(7, (6, 8)), alpha=1.0, zorder=5,
            ))
            # glow leve em ambos
            ax.add_patch(PathPatch(
                path, fill=False, edgecolor="#7090C8", linewidth=lw * 3,
                linestyle="-", alpha=0.08, zorder=3, capstyle="round",
            ))
            continue

        color, lw, ls, alpha = ROUTE_STYLE.get(
            label, ("#D4B574", 0.9, (0, (3, 3)), 0.70)
        )
        # glow embaixo — sutil
        ax.add_patch(PathPatch(
            path, fill=False, edgecolor=color, linewidth=lw * 2.8,
            linestyle="-", alpha=0.10, zorder=3, capstyle="round",
        ))
        # linha tracejada por cima — grossa
        ax.add_patch(PathPatch(
            path, fill=False, edgecolor=color, linewidth=lw * 1.5,
            linestyle=ls, alpha=alpha, zorder=5,
        ))

    # (D) Labels com leader lines — pra resolver amontoamento na Europa
    # mapeamento: label -> (label_x, label_y, ha)
    LABEL_POS = {
        "BRASIL · ORIGEM":  (-75, -25, "right"),
        "FRANÇA":           (-28, 60,  "right"),
        "SUÍÇA":            (-28, 50,  "right"),
        "SUÉCIA":           (-28, 70,  "right"),
        "ESPANHA · PRÓXIMO": (-25, 35, "right"),
        "OREGON · SONHO":   (-145, 35, "left"),
    }

    # (E) Pontos com glow (halo gold)
    for lon, lat, label, role in points:
        if role == "origin":
            for r, a in [(4.5, 0.10), (3.0, 0.18), (1.8, 0.35)]:
                ax.add_patch(Circle((lon, lat), r, color="#D4B574",
                                    alpha=a, zorder=5))
            ax.add_patch(Circle((lon, lat), 1.1, color="#F5E6C0", zorder=6))
        elif role == "next":
            for r, a in [(3.5, 0.10), (2.4, 0.18)]:
                ax.add_patch(Circle((lon, lat), r, color="#F0E0C0",
                                    alpha=a, zorder=5))
            ax.add_patch(Circle((lon, lat), 2.6, fill=False,
                                edgecolor="#F0E0C0", linewidth=1.0,
                                alpha=0.75, zorder=5))
            ax.add_patch(Circle((lon, lat), 0.9, color="#F0E0C0", zorder=6))
        elif role == "dream":
            for r, a in [(4.5, 0.10), (3.0, 0.18)]:
                ax.add_patch(Circle((lon, lat), r, color="#E8D4A8",
                                    alpha=a, zorder=5))
            ax.add_patch(Circle((lon, lat), 3.2, fill=False,
                                edgecolor="#E8D4A8", linewidth=0.8,
                                alpha=0.55, zorder=5))
            ax.add_patch(Circle((lon, lat), 1.8, fill=False,
                                edgecolor="#E8D4A8", linewidth=0.9,
                                alpha=0.85, zorder=5))
            ax.add_patch(Circle((lon, lat), 0.8, color="#F5E6C0", zorder=6))
        else:  # present
            for r, a in [(2.0, 0.12), (1.3, 0.22)]:
                ax.add_patch(Circle((lon, lat), r, color="#D4B574",
                                    alpha=a, zorder=5))
            ax.add_patch(Circle((lon, lat), 0.8, color="#D4B574", zorder=6))

        # leader line do ponto ate o label
        label_x, label_y, ha = LABEL_POS[label]
        ax.plot([lon, label_x], [lat, label_y],
                color="#8A7342", linewidth=0.5, alpha=0.7, zorder=4)

        # cor/tamanho do label
        if role == "origin":
            color, fs, fw = "#F5E6C0", 11, "bold"
        elif role == "next":
            color, fs, fw = "#F0E0C0", 10, "bold"
        elif role == "dream":
            color, fs, fw = "#F5E6C0", 10, "bold"
        else:
            color, fs, fw = "#D4B574", 9, "normal"

        # nudge do texto pra nao colar no fim do leader
        if ha == "right":
            tx = label_x - 1.5
        else:
            tx = label_x + 1.5
        ax.text(
            tx, label_y, label,
            color=color, fontsize=fs, fontfamily="Helvetica",
            fontweight=fw, ha=ha, va="center", zorder=7,
        )

    # (G) Vinheta nos cantos — 4 elipses pretas pra dar foco no centro
    for cx, cy in [(-200, 90), (200, 90), (-200, -75), (200, -75)]:
        ax.add_patch(Ellipse((cx, cy), width=200, height=120,
                              facecolor="#0A0A0A", edgecolor="none",
                              alpha=0.7, zorder=10))

    plt.tight_layout(pad=0)
    plt.savefig(out, facecolor="#0A0A0A", dpi=220, bbox_inches="tight", pad_inches=0.1)
    plt.close(fig)
    print(f"  mapa ok: {out}")
    return out

WORLD_MAP = build_world_map()

# -----------------------------------------------------------------------------
# Gradientes
# -----------------------------------------------------------------------------
def make_gradient_png(name, w, h, direction="vertical", start_alpha=255, end_alpha=0):
    out = f"{IMG_DIR}/grad-{name}.png"
    img = Image.new("RGBA", (w, h), (0, 0, 0, 0))
    px = img.load()
    if direction == "vertical":
        for y in range(h):
            t = y / max(h - 1, 1)
            a = int(start_alpha + (end_alpha - start_alpha) * t)
            for x in range(w):
                px[x, y] = (0, 0, 0, a)
    else:
        for x in range(w):
            t = x / max(w - 1, 1)
            a = int(start_alpha + (end_alpha - start_alpha) * t)
            for y in range(h):
                px[x, y] = (0, 0, 0, a)
    img.save(out, "PNG", optimize=True)
    return out

GRAD_BOTTOM = make_gradient_png("bottom", 1600, 900, "vertical", 255, 0)
GRAD_TOP = make_gradient_png("top", 1600, 900, "vertical", 0, 255)

# -----------------------------------------------------------------------------
# Helpers de slide
# -----------------------------------------------------------------------------
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = prs.slide_width
SH = prs.slide_height
BLANK = prs.slide_layouts[6]

def add_bg(slide, color=BLACK):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.shadow.inherit = False
    sp = bg.shadow
    return bg

def add_image_fullbleed(slide, path):
    with Image.open(path) as im:
        iw, ih = im.size
    slide_ratio = SW / SH
    img_ratio = iw / ih
    if img_ratio > slide_ratio:
        new_h = SH
        new_w = int(new_h * img_ratio)
        left = (SW - new_w) // 2
        top = 0
    else:
        new_w = SW
        new_h = int(new_w / img_ratio)
        left = 0
        top = (SH - new_h) // 2
    slide.shapes.add_picture(path, left, top, width=new_w, height=new_h)

def add_overlay(slide, alpha_pct):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    rect.line.fill.background()
    rect.fill.solid()
    rect.fill.fore_color.rgb = BLACK
    solid_fill = rect.fill._xPr.find(qn("a:solidFill"))
    if solid_fill is None:
        solid_fill = rect.fill._xPr.find(".//" + qn("a:solidFill"))
    srgb = solid_fill.find(qn("a:srgbClr"))
    alpha_val = int(alpha_pct * 1000)
    alpha_el = etree.SubElement(srgb, qn("a:alpha"))
    alpha_el.set("val", str(alpha_val))
    return rect

def add_gradient_overlay(slide, path, height_frac=1.0, anchor="bottom"):
    h = int(SH * height_frac)
    if anchor == "bottom":
        top = SH - h
    else:
        top = 0
    slide.shapes.add_picture(path, 0, top, width=SW, height=h)

def add_text(slide, left, top, width, height, text, *,
             font=F_SERIF, size=24, color=CREAM, bold=False, italic=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, letter_spacing=None,
             line_spacing=1.15):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    lines = text.split("<br>")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        f = r.font
        f.name = font
        f.size = Pt(size)
        f.bold = bold
        f.italic = italic
        f.color.rgb = color
        if letter_spacing is not None:
            rPr = r._r.get_or_add_rPr()
            rPr.set("spc", str(int(letter_spacing * 100)))
    return tb

def add_hline(slide, left, top, width, color=GOLD, weight=0.75):
    line = slide.shapes.add_connector(1, left, top, left + width, top)
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    return line

def add_vline(slide, left, top, height, color=GOLD, weight=0.5):
    line = slide.shapes.add_connector(1, left, top, left, top + height)
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    return line

def add_dot(slide, cx, cy, r_emu, color=GOLD):
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                 cx - r_emu, cy - r_emu, r_emu * 2, r_emu * 2)
    dot.line.fill.background()
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    return dot

def crop_cover(src_path, target_w_emu, target_h_emu, out_name):
    out = f"{IMG_DIR}/crop-{out_name}.jpg"
    target_ratio = target_w_emu / target_h_emu
    with Image.open(src_path) as im:
        iw, ih = im.size
        ir = iw / ih
        if ir > target_ratio:
            new_w = int(ih * target_ratio)
            left = (iw - new_w) // 2
            box = (left, 0, left + new_w, ih)
        else:
            new_h = int(iw / target_ratio)
            top = (ih - new_h) // 2
            box = (0, top, iw, top + new_h)
        im2 = im.crop(box)
        im2.save(out, "JPEG", quality=92, optimize=True)
    return out

# -----------------------------------------------------------------------------
# SLIDE 1 — CAPA (PT) — usa cherries-yellow
# -----------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_bg(s, BLACK)
# usa cherries-yellow cropado em fullbleed
cover_crop = crop_cover(JPGS["cherries-yellow"], SW, SH, "s1-cover-pt")
s.shapes.add_picture(cover_crop, 0, 0, width=SW, height=SH)
add_overlay(s, 55)
add_gradient_overlay(s, GRAD_BOTTOM, height_frac=0.7, anchor="bottom")

WM_W = Inches(2.0)
WM_H = Inches(0.5)
s.shapes.add_picture(WORDMARK_GOLD,
                     (SW - WM_W) // 2, Inches(0.45),
                     width=WM_W, height=WM_H)

add_text(s, Inches(0), Inches(2.4), SW, Inches(0.4),
         "VALIM FARMS  ·  SÃO PAULO  ·  BRASIL",
         font=F_SANS, size=11, color=GOLD_LIGHT, align=PP_ALIGN.CENTER,
         letter_spacing=5, bold=False)

# Headline removido — apresentador escreve a frase ao vivo / na hora
# (mantemos somente wordmark, eyebrow, linha dourada e atribuicao)

add_hline(s, Inches(5.5), Inches(6.35), Inches(2.333), GOLD, 0.75)

add_text(s, Inches(0), Inches(6.55), SW, Inches(0.4),
         "WILSON LUIZ VALIM ZERBINATTI  ·  TERCEIRA GERAÇÃO  ·  1897",
         font=F_SANS, size=10, color=GOLD, align=PP_ALIGN.CENTER,
         letter_spacing=4)

# -----------------------------------------------------------------------------
# SLIDE 2 — LINHAGEM (PT)
# -----------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_bg(s, BLACK)

img_w = Inches(6.0)
img_h = SH

crop1 = crop_cover(JPGS["flowers-sunset"], img_w, img_h, "s2-left")
s.shapes.add_picture(crop1, 0, 0, width=img_w, height=img_h)

ov = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, img_w, SH)
ov.line.fill.background()
ov.fill.solid()
ov.fill.fore_color.rgb = BLACK
srgb = ov.fill._xPr.find(".//" + qn("a:srgbClr"))
alpha_el = etree.SubElement(srgb, qn("a:alpha"))
alpha_el.set("val", "25000")

RX = Inches(6.7)
RW = Inches(6.0)

add_text(s, RX, Inches(1.2), RW, Inches(0.4),
         "— LINHAGEM",
         font=F_SANS, size=11, color=GOLD, letter_spacing=5)

add_text(s, RX, Inches(1.7), RW, Inches(2.7),
         "Três gerações.<br>Uma terra.<br>Um café.",
         font=F_SERIF, size=54, color=CREAM, italic=True, line_spacing=1.1)

add_hline(s, RX, Inches(4.55), Inches(0.7), GOLD, 0.75)

body = ("Família italiana chegou em São Paulo em 1897. Plantou. Esperou. Aprendeu. "
        "Hoje sou eu, Wilson Luiz Valim Zerbinatti, o terceiro a tocar o cafezal. "
        "O sobrenome Zerbinatti virou marca — mas continua sendo o que era: "
        "pessoas que escutam a terra.")
add_text(s, RX, Inches(4.85), RW, Inches(1.6),
         body, font=F_SANS, size=13, color=CREAM, line_spacing=1.4)

add_hline(s, RX, Inches(6.45), Inches(5.5), GOLD, 0.5)
add_text(s, RX, Inches(6.6), RW, Inches(0.5),
         "129 ANOS  ·  TRÊS GERAÇÕES  ·  UMA ÚNICA ORIGEM",
         font=F_SANS, size=11, color=GOLD_LIGHT, letter_spacing=4, bold=True)

# -----------------------------------------------------------------------------
# SLIDE 3 — A TERRA (PT)
# -----------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_bg(s, BLACK)

right_w = Inches(8.0)
right_h = SH
crop2 = crop_cover(JPGS["cherries-yellow"], right_w, right_h, "s3-right")
s.shapes.add_picture(crop2, SW - right_w, 0, width=right_w, height=right_h)

grad_lr = make_gradient_png("lr-s3", 1200, 900, "horizontal", 255, 0)
s.shapes.add_picture(grad_lr, SW - right_w - Inches(1.5), 0,
                     width=Inches(3.5), height=SH)

ov = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, SW - right_w, 0, right_w, SH)
ov.line.fill.background()
ov.fill.solid()
ov.fill.fore_color.rgb = BLACK
srgb = ov.fill._xPr.find(".//" + qn("a:srgbClr"))
alpha_el = etree.SubElement(srgb, qn("a:alpha"))
alpha_el.set("val", "30000")

LX = Inches(0.7)
LW = Inches(5.5)

add_text(s, LX, Inches(1.0), LW, Inches(0.4),
         "— A TERRA",
         font=F_SANS, size=11, color=GOLD, letter_spacing=5)

add_text(s, LX, Inches(1.5), LW, Inches(2.2),
         "Onde tudo<br>acontece sozinho.",
         font=F_SERIF, size=52, color=CREAM, italic=True, line_spacing=1.05)

table_data = [
    ("LOCALIZAÇÃO", "Valim Farms, SP"),
    ("ALTITUDE", "640 – 760 m"),
    ("SOLO", "Vulcânico"),
    ("VARIEDADES", "Arara · Paraíso 2"),
    ("PROCESSO", "Natural, ao sol"),
    ("COLHEITA", "Manual, cereja madura"),
]
T_TOP = Inches(4.0)
ROW_H = Inches(0.45)
COL1_W = Inches(2.0)
COL2_W = Inches(3.5)
for i, (k, v) in enumerate(table_data):
    y = T_TOP + ROW_H * i
    add_hline(s, LX, y, COL1_W + COL2_W, GOLD_DARK, 0.5)
    add_text(s, LX, y + Inches(0.10), COL1_W, ROW_H,
             k, font=F_SANS, size=9.5, color=GOLD,
             letter_spacing=4, bold=True)
    add_text(s, LX + COL1_W, y + Inches(0.08), COL2_W, ROW_H,
             v, font=F_SERIF, size=15, color=CREAM, italic=True)
add_hline(s, LX, T_TOP + ROW_H * len(table_data), COL1_W + COL2_W,
          GOLD_DARK, 0.5)

# Tira de selos no rodape do lado esquerdo (credibilidade institucional)
SELOS = [
    f"{ROOT}/public/assets/selo-cup.png",
    f"{ROOT}/public/assets/selo-sca-90-50.png",
    f"{ROOT}/public/assets/selo-organico-fazenda.png",
    f"{ROOT}/public/assets/selo-brasil.png",
]
SELO_SIZE = Inches(0.55)
SELO_AREA_LEFT = Inches(0.7)
SELO_AREA_RIGHT = Inches(6.2)
SELO_AREA_W = SELO_AREA_RIGHT - SELO_AREA_LEFT
SELO_Y = Inches(6.85)
# eyebrow acima dos selos
add_text(s, SELO_AREA_LEFT, Inches(6.55), SELO_AREA_W, Inches(0.3),
         "— CERTIFICAÇÕES",
         font=F_SANS, size=9, color=GOLD_DARK, letter_spacing=4, bold=True)
# espacamento uniforme entre centros dos selos
n_selos = len(SELOS)
slot_w = SELO_AREA_W // n_selos
for idx, selo_path in enumerate(SELOS):
    center_x = SELO_AREA_LEFT + slot_w * idx + slot_w // 2
    selo_left = center_x - SELO_SIZE // 2
    s.shapes.add_picture(selo_path, selo_left, SELO_Y,
                         width=SELO_SIZE, height=SELO_SIZE)

# -----------------------------------------------------------------------------
# SLIDE 4 — BRAND PORTFOLIO / VALIM FARMS · TERROIR (PT)
# -----------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_bg(s, BLACK)

# TOPO: eyebrow esquerdo + italic ouro direito
add_text(s, Inches(0.7), Inches(0.55), Inches(8), Inches(0.4),
         "ZERBINATTI  ·  BRAND PORTFOLIO",
         font=F_SANS, size=9, color=GOLD_DARK, letter_spacing=4, bold=True)

add_text(s, Inches(9.5), Inches(0.55), Inches(3.2), Inches(0.4),
         "Famiglia · since 1897",
         font=F_SERIF, size=12, color=GOLD, italic=True,
         align=PP_ALIGN.RIGHT)

# HEADLINE + BODY
add_text(s, Inches(0.7), Inches(1.3), Inches(8), Inches(0.4),
         "— VALIM FARMS  ·  TERROIR",
         font=F_SANS, size=10, color=GOLD, letter_spacing=5, bold=True)

add_text(s, Inches(0.7), Inches(1.7), Inches(6.5), Inches(1.5),
         "Uma fazenda.<br>Um clima inteiro.",
         font=F_SERIF, size=44, color=CREAM, italic=True, line_spacing=1.05)

s4_body_pt = ("A Valim Farms fica em terraços naturais entre 640 e 760 metros, "
              "com solo vulcânico, invernos secos e uma janela de maturação longa e lenta. "
              "Duas variedades arábica — Arara e Paraíso 2 — plantadas à mão e colhidas "
              "só na cereja madura.")
add_text(s, Inches(7.5), Inches(1.85), Inches(5.1), Inches(1.3),
         s4_body_pt, font=F_SANS, size=11, color=CREAM, line_spacing=1.45)

# LINHA DIVISORA
add_hline(s, Inches(0.7), Inches(3.4), Inches(11.93), GOLD_DARK, 0.5)

# GRID FOTOS
# Coluna A — foto principal grande
s4_w_a = Inches(4.2)
s4_h_a = Inches(3.2)
s4_crop_a = crop_cover(JPGS["cherries-yellow"], s4_w_a, s4_h_a, "s4-photo-a")
s.shapes.add_picture(s4_crop_a, Inches(0.7), Inches(3.7),
                     width=s4_w_a, height=s4_h_a)

# Coluna B — duas fotos empilhadas
s4_w_b = Inches(3.5)
s4_h_b = Inches(1.55)
s4_crop_b1 = crop_cover(JPGS["young-plants"], s4_w_b, s4_h_b, "s4-photo-b1")
s.shapes.add_picture(s4_crop_b1, Inches(5.0), Inches(3.7),
                     width=s4_w_b, height=s4_h_b)
s4_crop_b2 = crop_cover(JPGS["drying-yard-tree"], s4_w_b, s4_h_b, "s4-photo-b2")
s.shapes.add_picture(s4_crop_b2, Inches(5.0), Inches(5.35),
                     width=s4_w_b, height=s4_h_b)

# Coluna C — CARD TERROIR
card = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                          Inches(8.7), Inches(3.7),
                          Inches(3.93), Inches(3.2))
card.fill.solid()
card.fill.fore_color.rgb = BLACK
card.line.color.rgb = GOLD_DARK
card.line.width = Pt(0.5)

# Headline do card
add_text(s, Inches(8.95), Inches(3.95), Inches(3.4), Inches(0.5),
         "O terroir.",
         font=F_SERIF, size=22, color=CREAM, italic=True)

# Body do card
s4_card_body_pt = ("Solo vulcânico, encostas em terraço, manhãs frias, invernos secos. "
                   "A terra faz o trabalho — a gente só não atrapalha.")
add_text(s, Inches(8.95), Inches(4.55), Inches(3.4), Inches(1.4),
         s4_card_body_pt, font=F_SANS, size=10, color=CREAM, line_spacing=1.45)

# Linha divisora dentro do card
add_hline(s, Inches(8.95), Inches(5.85), Inches(3.4), GOLD_DARK, 0.5)

# Stats grid 2x2
# Esquerda top: 760m / ALTITUDE MAX
add_text(s, Inches(8.95), Inches(6.05), Inches(1.7), Inches(0.4),
         "760m", font=F_SERIF, size=20, color=GOLD, italic=True)
add_text(s, Inches(8.95), Inches(6.40), Inches(1.7), Inches(0.3),
         "ALTITUDE MAX", font=F_SANS, size=8, color=GOLD_DARK,
         letter_spacing=4, bold=True)

# Direita top: 02 / VARIEDADES
add_text(s, Inches(10.75), Inches(6.05), Inches(1.7), Inches(0.4),
         "02", font=F_SERIF, size=20, color=GOLD, italic=True)
add_text(s, Inches(10.75), Inches(6.40), Inches(1.7), Inches(0.3),
         "VARIEDADES", font=F_SANS, size=8, color=GOLD_DARK,
         letter_spacing=4, bold=True)

# Esquerda bottom: 100% / ARÁBICA
add_text(s, Inches(8.95), Inches(6.55), Inches(1.7), Inches(0.4),
         "100%", font=F_SERIF, size=20, color=GOLD, italic=True)
add_text(s, Inches(8.95), Inches(6.90), Inches(1.7), Inches(0.3),
         "ARÁBICA", font=F_SANS, size=8, color=GOLD_DARK,
         letter_spacing=4, bold=True)

# Direita bottom: Mai-Set / COLHEITA
add_text(s, Inches(10.75), Inches(6.55), Inches(1.7), Inches(0.4),
         "Mai-Set", font=F_SERIF, size=20, color=GOLD, italic=True)
add_text(s, Inches(10.75), Inches(6.90), Inches(1.7), Inches(0.3),
         "COLHEITA", font=F_SANS, size=8, color=GOLD_DARK,
         letter_spacing=4, bold=True)

# -----------------------------------------------------------------------------
# SLIDE 5 — O MUNDO (PT)
# -----------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_bg(s, BLACK)

add_text(s, Inches(0.7), Inches(0.7), Inches(12), Inches(0.4),
         "— O CAFÉ JÁ ATRAVESSA OCEANOS",
         font=F_SANS, size=11, color=GOLD, letter_spacing=5, bold=True)

add_text(s, Inches(0.7), Inches(1.15), Inches(12), Inches(1.6),
         "Hoje, três bandeiras.<br>O sonho é Oregon.",
         font=F_SERIF, size=42, color=CREAM, italic=True, line_spacing=1.1)

map_w = Inches(11.5)
map_h = Inches(4.2)
map_left = (SW - map_w) // 2
map_top = Inches(2.8)
s.shapes.add_picture(WORLD_MAP, map_left, map_top, width=map_w, height=map_h)

add_text(s, Inches(0), Inches(7.0), SW, Inches(0.4),
         "PRESENÇA  ·  EXPANSÃO  ·  2026",
         font=F_SANS, size=11, color=GOLD_LIGHT,
         letter_spacing=5, align=PP_ALIGN.CENTER, bold=True)

# -----------------------------------------------------------------------------
# SLIDE 6 — OS NUMEROS (PT)
# -----------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_bg(s, BLACK)

def make_radial(name, w=1600, h=900, color=(201, 169, 97), max_alpha=70):
    out = f"{IMG_DIR}/radial-{name}.png"
    img = Image.new("RGBA", (w, h), (0, 0, 0, 0))
    px = img.load()
    cx, cy = w / 2, h / 2
    rmax = ((w / 2) ** 2 + (h / 2) ** 2) ** 0.5 * 0.7
    for y in range(h):
        for x in range(w):
            d = ((x - cx) ** 2 + (y - cy) ** 2) ** 0.5
            t = min(d / rmax, 1.0)
            a = int(max_alpha * (1 - t) ** 2)
            px[x, y] = (color[0], color[1], color[2], a)
    img.save(out, "PNG", optimize=True)
    return out

radial = make_radial("s5", 1200, 700, (201, 169, 97), 55)
s.shapes.add_picture(radial, 0, 0, width=SW, height=SH)

add_text(s, Inches(0.7), Inches(0.8), Inches(12), Inches(0.4),
         "— A PROVA",
         font=F_SANS, size=11, color=GOLD, letter_spacing=5, bold=True)

add_text(s, Inches(0.7), Inches(1.3), Inches(12), Inches(1.6),
         "Os números<br>que importam.",
         font=F_SERIF, size=46, color=CREAM, italic=True, line_spacing=1.1)

# HERO NUMBER (esquerda) — placeholder gigante de sacas/ano
HERO_LEFT = Inches(0.7)
HERO_TOP = Inches(3.5)
HERO_W = Inches(6.5)
add_text(s, HERO_LEFT, HERO_TOP, HERO_W, Inches(2.4),
         "[ ___ ]", font=F_SERIF, size=160, color=GOLD_DARK,
         italic=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
         line_spacing=1.0)
add_text(s, HERO_LEFT, HERO_TOP + Inches(2.5), HERO_W, Inches(0.4),
         "SACAS / ANO", font=F_SANS, size=12, color=CREAM,
         letter_spacing=5, align=PP_ALIGN.CENTER, bold=True)
add_text(s, HERO_LEFT, HERO_TOP + Inches(2.95), HERO_W, Inches(0.4),
         "Produção 2026", font=F_SERIF, size=13, color=GOLD_LIGHT,
         italic=True, align=PP_ALIGN.CENTER)

# linha vertical divisoria
add_vline(s, Inches(7.2), Inches(3.7), Inches(3.2), GOLD_DARK, 0.5)

# 3 STATS MENORES (direita), empilhados
STATS_X = Inches(7.5)
STATS_W = Inches(5.3)
STATS = [
    ("20 CONTAINERS",  "EMBARQUE ESPANHA · 2026",  "value"),
    ("[ __ ]",         "PONTUAÇÃO SCAA",            "placeholder"),
    ("100%",           "MICROLOTE RASTREÁVEL",      "value"),
]
SUBLABELS = [
    "Embarque programado",
    "Cup rating",
    "Origem garantida",
]
STATS_TOP = Inches(3.7)
STATS_BLOCK_H = Inches(1.07)  # 3.2 / 3
for i, ((num, label, kind), sub) in enumerate(zip(STATS, SUBLABELS)):
    y = STATS_TOP + STATS_BLOCK_H * i
    color_num = GOLD_DARK if kind == "placeholder" else GOLD
    add_text(s, STATS_X, y, STATS_W, Inches(0.6),
             num, font=F_SERIF, size=38, color=color_num,
             italic=True, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             line_spacing=1.0)
    add_text(s, STATS_X, y + Inches(0.6), STATS_W, Inches(0.3),
             label, font=F_SANS, size=10, color=GOLD,
             letter_spacing=4, bold=True, align=PP_ALIGN.LEFT)
    add_text(s, STATS_X, y + Inches(0.88), STATS_W, Inches(0.3),
             sub, font=F_SERIF, size=11, color=GOLD_LIGHT,
             italic=True, align=PP_ALIGN.LEFT)

# -----------------------------------------------------------------------------
# SLIDE 7 — ESTADOS UNIDOS (PT)
# -----------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_bg(s, BLACK)

add_image_fullbleed(s, JPGS["drying-yard-tree"])
add_overlay(s, 72)

add_gradient_overlay(s, GRAD_BOTTOM, height_frac=1.0, anchor="bottom")
ov = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
ov.line.fill.background()
ov.fill.solid()
ov.fill.fore_color.rgb = BLACK
srgb = ov.fill._xPr.find(".//" + qn("a:srgbClr"))
alpha_el = etree.SubElement(srgb, qn("a:alpha"))
alpha_el.set("val", "30000")

add_text(s, Inches(0), Inches(0.7), SW, Inches(0.4),
         "— MEU PRÓXIMO DESTINO",
         font=F_SANS, size=11, color=GOLD, letter_spacing=5,
         align=PP_ALIGN.CENTER, bold=True)

add_text(s, Inches(0.3), Inches(1.7), Inches(12.7), Inches(3.3),
         "Estados Unidos.",
         font=F_SERIF, size=180, color=GOLD, italic=True,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)

add_hline(s, Inches(5.0), Inches(5.5), Inches(3.333), GOLD, 1.0)

add_text(s, Inches(1.5), Inches(5.7), Inches(10.333), Inches(1.2),
         "Deixaria de exportar para a Europa<br>para enfim chegar aqui.",
         font=F_SERIF, size=28, color=CREAM, italic=True,
         align=PP_ALIGN.CENTER, line_spacing=1.25)

add_text(s, Inches(0), Inches(7.05), SW, Inches(0.35),
         "WILSON LUIZ VALIM ZERBINATTI  ·  2026",
         font=F_SANS, size=10, color=GOLD_LIGHT, letter_spacing=4,
         align=PP_ALIGN.CENTER, bold=True)

# -----------------------------------------------------------------------------
# SLIDE 8 — ENCERRAMENTO (PT)
# -----------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_bg(s, BLACK)
add_image_fullbleed(s, JPGS["peneirar"])
add_overlay(s, 78)

add_text(s, Inches(1.0), Inches(0.9), Inches(11.333), Inches(2.3),
         "\u201CA terra ensina, o sol decide,<br>"
         "a gente apenas observa.\u201D",
         font=F_SERIF, size=36, color=CREAM, italic=True,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.3)

orn_y = Inches(3.45)
add_hline(s, Inches(4.0), orn_y, Inches(2.3), GOLD, 0.75)
add_hline(s, Inches(7.0), orn_y, Inches(2.3), GOLD, 0.75)
diamond_r = Inches(0.08)
diamond = s.shapes.add_shape(MSO_SHAPE.DIAMOND,
                              Inches(6.667) - diamond_r,
                              orn_y - diamond_r,
                              diamond_r * 2, diamond_r * 2)
diamond.line.color.rgb = GOLD
diamond.line.width = Pt(0.5)
diamond.fill.solid()
diamond.fill.fore_color.rgb = GOLD

add_text(s, Inches(0), Inches(3.7), SW, Inches(0.4),
         "— WILSON LUIZ VALIM ZERBINATTI  ·  TERCEIRA GERAÇÃO",
         font=F_SANS, size=9, color=GOLD, letter_spacing=5,
         align=PP_ALIGN.CENTER, bold=True)

# QR code centralizado (medalhao redondo)
QR_SIZE = Inches(2.1)
qr_x = (SW - QR_SIZE) / 2
qr_y = Inches(4.4)
s.shapes.add_picture(QR_PATH, qr_x, qr_y, width=QR_SIZE, height=QR_SIZE)

# Rodape: Instagram + site abaixo do QR
add_text(s, Inches(0), Inches(6.75), SW, Inches(0.4),
         "@ZERBINATTICOFFEE  ·  ZERBINATTICOFFEE.COM",
         font=F_SANS, size=12, color=GOLD_LIGHT, letter_spacing=3,
         align=PP_ALIGN.CENTER, bold=True)

WM2_W = Inches(1.2)
WM2_H = Inches(0.3)
s.shapes.add_picture(WORDMARK_GOLD,
                     SW - WM2_W - Inches(0.4),
                     SH - WM2_H - Inches(0.25),
                     width=WM2_W, height=WM2_H)

# -----------------------------------------------------------------------------
# SALVAR
# -----------------------------------------------------------------------------
prs.save(OUT_PATH)
size_kb = os.path.getsize(OUT_PATH) / 1024
print(f"\nOK -> {OUT_PATH}")
print(f"slides: {len(prs.slides)}  size: {size_kb:.1f} KB")
