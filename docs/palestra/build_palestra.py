"""
Gerador da palestra Zerbinatti (5min) — slide deck editorial premium.
Saida: /Users/fabricio/dev/cafe/docs/palestra/zerbinatti-palestra-5min.pptx
"""
import os
import json
import urllib.request
from io import BytesIO
from PIL import Image, ImageFilter

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.patches import Circle, Ellipse, FancyArrowPatch, PathPatch
from matplotlib.path import Path as MplPath
import matplotlib.patheffects as path_effects
import numpy as np

# -----------------------------------------------------------------------------
# Constantes de identidade visual
# -----------------------------------------------------------------------------
BLACK = RGBColor(0x0A, 0x0A, 0x0A)
CREAM = RGBColor(0xEC, 0xEC, 0xEC)
GOLD = RGBColor(0xC9, 0xA9, 0x61)
GOLD_LIGHT = RGBColor(0xDD, 0xC8, 0xA0)
GOLD_PALE = RGBColor(0xF0, 0xE0, 0xC0)
GOLD_DARK = RGBColor(0x8A, 0x73, 0x42)

# Fontes (sistema macOS)
F_SERIF = "Cormorant Garamond"
F_SERIF_FB = "Garamond"
F_SANS = "Inter"
F_SANS_FB = "Helvetica Neue"

# Caminhos
ROOT = "/Users/fabricio/dev/cafe"
OUT_DIR = f"{ROOT}/docs/palestra"
IMG_DIR = f"{OUT_DIR}/img"
OUT_PATH = f"{OUT_DIR}/zerbinatti-palestra-5min.pptx"

os.makedirs(IMG_DIR, exist_ok=True)

# -----------------------------------------------------------------------------
# Converter webp -> jpg (python-pptx nao le webp)
# -----------------------------------------------------------------------------
WEBP_FILES = {
    "cherries-yellow": f"{ROOT}/public/images/farm/cherries-yellow.webp",
    "drying-yard-tree": f"{ROOT}/public/images/farm/drying-yard-tree.webp",
    "flowers-sunset": f"{ROOT}/public/images/farm/flowers-sunset.webp",
    "young-plants": f"{ROOT}/public/images/farm/young-plants.webp",
    "seedlings-nursery": f"{ROOT}/public/images/farm/seedlings-nursery.webp",
    "packages-row": f"{ROOT}/public/images/farm/packages-row.webp",
    "peneirar": f"{ROOT}/public/assets/galeria/peneirar.webp",
}

JPGS = {}
for name, src in WEBP_FILES.items():
    dst = f"{IMG_DIR}/{name}.jpg"
    if not os.path.exists(dst):
        im = Image.open(src).convert("RGB")
        im.save(dst, "JPEG", quality=92, optimize=True)
    JPGS[name] = dst
    print(f"  jpg ok: {dst}")

# Hero da capa (foto do site, JPG ja em RGB) — upscale via LANCZOS para 2400x2400
HERO_SRC = f"{ROOT}/public/assets/hero-bg.jpg"
HERO_JPG = f"{IMG_DIR}/hero-bg-upscaled.jpg"
if not os.path.exists(HERO_JPG):
    with Image.open(HERO_SRC) as im:
        im = im.convert("RGB")
        # upscale para 2400 mantendo proporcao quadrada (743x723 ~ quadrada)
        target = 2400
        iw, ih = im.size
        ratio = target / max(iw, ih)
        new_size = (int(iw * ratio), int(ih * ratio))
        im_up = im.resize(new_size, Image.LANCZOS)
        # leve sharpen para compensar upscale
        im_up = im_up.filter(ImageFilter.UnsharpMask(radius=1.2, percent=120, threshold=2))
        im_up.save(HERO_JPG, "JPEG", quality=94, optimize=True)
print(f"  hero ok: {HERO_JPG}")

WORDMARK_GOLD = f"{ROOT}/public/assets/zerbinatti-wordmark-gold.png"
WORDMARK_CREAM = f"{ROOT}/public/assets/zerbinatti-wordmark-cream.png"
QR_PATH = f"{IMG_DIR}/qr-zerbinatti.png"

# -----------------------------------------------------------------------------
# Geracao do mapa-mundi estilizado em ouro/preto para slide 4
# -----------------------------------------------------------------------------
GEO_URL = "https://raw.githubusercontent.com/nvkelso/natural-earth-vector/master/geojson/ne_110m_admin_0_countries.geojson"
GEO_PATH = f"{IMG_DIR}/ne_110m_countries.geojson"


def build_world_map():
    """Mapa-mundi minimalista em ouro/preto usando contornos reais (Natural Earth 110m).
    Marca: Brasil (origem), Franca/Suica/Suecia (presenca), Espanha (proximo), Oregon (sonho).
    Cada destino recebe uma rota tracejada saindo do Brasil.
    """
    out = f"{IMG_DIR}/world-map.png"

    # Cache do GeoJSON Natural Earth (low-res, ~200KB)
    if not os.path.exists(GEO_PATH):
        print(f"  baixando geojson: {GEO_URL}")
        urllib.request.urlretrieve(GEO_URL, GEO_PATH)
    with open(GEO_PATH) as f:
        world = json.load(f)

    fig, ax = plt.subplots(figsize=(16, 8), facecolor="#0A0A0A")
    ax.set_facecolor("#0A0A0A")
    ax.set_xlim(-180, 180)
    ax.set_ylim(-58, 80)
    ax.set_aspect("equal")
    ax.axis("off")

    GOLD_HEX = "#D4B574"     # borda mais clara/luminosa
    GOLD_FILL = "#4A3A1F"    # fill mais quente

    # (A) Atmosfera de fundo (oceano) — gradiente radial sutil
    for r, alpha in [(220, 0.025), (160, 0.04), (100, 0.05)]:
        ax.add_patch(Ellipse((0, 10), width=r * 2, height=r * 1.1,
                              facecolor="#C9A961", edgecolor="none",
                              alpha=alpha, zorder=0))

    # (C) Equador + tropicos (linhas pontilhadas discretas)
    for y_lat, a_line in [(0, 0.25), (23.5, 0.12), (-23.5, 0.12)]:
        ax.axhline(y=y_lat, color="#C9A961", linewidth=0.4,
                   linestyle=(0, (1, 3)), alpha=a_line, zorder=1)

    # (B) Países com profundidade — borda mais grossa e cores mais quentes
    for feat in world["features"]:
        if feat.get("properties", {}).get("NAME") == "Antarctica":
            continue
        geom = feat.get("geometry")
        if geom is None:
            continue
        gtype = geom["type"]
        if gtype == "MultiPolygon":
            polys = geom["coordinates"]
        elif gtype == "Polygon":
            polys = [geom["coordinates"]]
        else:
            continue
        for poly in polys:
            ring = poly[0]
            patch = mpatches.Polygon(
                ring, closed=True,
                facecolor=GOLD_FILL, edgecolor=GOLD_HEX,
                linewidth=0.9, alpha=0.95, joinstyle="round",
                zorder=2,
            )
            ax.add_patch(patch)

    points = [
        (-47, -15, "BRAZIL · ORIGIN", "origin"),
        (2,   47, "FRANCE",           "present"),
        (8,   46, "SWITZERLAND",      "present"),
        (16,  60, "SWEDEN",           "present"),
        (-4,  40, "SPAIN · NEXT",     "next"),
        (-122, 45, "OREGON · DREAM",  "dream"),
    ]

    # (F) Rotas com glow — desenha 2x: glow grosso + linha tracejada
    # Cada destino tem altura de arco diferente -> efeito arco-iris empilhado
    ARC_HEIGHT = {
        "SPAIN · NEXT":    0.10,
        "FRANCE":          0.32,
        "SWITZERLAND":     0.56,
        "SWEDEN":          0.82,
        "OREGON · DREAM":  0.22,
    }
    # cor distinta por destino — arco-iris empilhado
    ROUTE_STYLE = {
        "SPAIN · NEXT":    ("#FF4FA3", 1.8, (0, (6, 4)),   1.00),  # rosa
        "FRANCE":          ("#FFD93C", 1.5, (0, (5, 4)),   1.00),  # amarelo
        "SWITZERLAND":     ("#3CD982", 1.5, (0, (4, 4)),   1.00),  # verde
        "SWEDEN":          ("#A878FF", 1.5, (0, (3, 4)),   1.00),  # roxo claro
    }
    OREGON_LABEL = "OREGON · DREAM"
    origin = (-47, -15)
    for lon, lat, label, role in points:
        if role == "origin":
            continue
        mid_x = (origin[0] + lon) / 2
        h = ARC_HEIGHT.get(label, 0.18)
        mid_y = (origin[1] + lat) / 2 + (abs(lon - origin[0]) * h)
        verts = [origin, (mid_x, mid_y), (lon, lat)]
        codes = [MplPath.MOVETO, MplPath.CURVE3, MplPath.CURVE3]
        path = MplPath(verts, codes)

        if label == OREGON_LABEL:
            lw = 1.6
            ax.add_patch(PathPatch(
                path, fill=False, edgecolor="#1E5BD8", linewidth=lw * 1.4,
                linestyle=(0, (6, 8)), alpha=1.0, zorder=5,
            ))
            ax.add_patch(PathPatch(
                path, fill=False, edgecolor="#E63946", linewidth=lw * 1.4,
                linestyle=(7, (6, 8)), alpha=1.0, zorder=5,
            ))
            ax.add_patch(PathPatch(
                path, fill=False, edgecolor="#7090C8", linewidth=lw * 3,
                linestyle="-", alpha=0.08, zorder=3, capstyle="round",
            ))
            continue

        color, lw, ls, alpha = ROUTE_STYLE.get(
            label, ("#D4B574", 0.9, (0, (3, 3)), 0.70)
        )
        ax.add_patch(PathPatch(
            path, fill=False, edgecolor=color, linewidth=lw * 2.8,
            linestyle="-", alpha=0.10, zorder=3, capstyle="round",
        ))
        ax.add_patch(PathPatch(
            path, fill=False, edgecolor=color, linewidth=lw * 1.5,
            linestyle=ls, alpha=alpha, zorder=5,
        ))

    # (D) Labels com leader lines — pra resolver amontoamento na Europa
    LABEL_POS = {
        "BRAZIL · ORIGIN":  (-75, -25, "right"),
        "FRANCE":           (-28, 60,  "right"),
        "SWITZERLAND":      (-28, 50,  "right"),
        "SWEDEN":           (-28, 70,  "right"),
        "SPAIN · NEXT":     (-25, 35,  "right"),
        "OREGON · DREAM":   (-145, 35, "left"),
    }

    # (E) Pontos com glow (halo gold)
    for lon, lat, label, role in points:
        if role == "origin":
            for r, a in [(4.5, 0.10), (3.0, 0.18), (1.8, 0.35)]:
                ax.add_patch(Circle((lon, lat), r, color="#D4B574",
                                    alpha=a, zorder=5))
            ax.add_patch(Circle((lon, lat), 1.1, color="#F5E6C0", zorder=6))
        elif role == "next":
            for r, a in [(3.5, 0.10), (2.4, 0.18)]:
                ax.add_patch(Circle((lon, lat), r, color="#F0E0C0",
                                    alpha=a, zorder=5))
            ax.add_patch(Circle((lon, lat), 2.6, fill=False,
                                edgecolor="#F0E0C0", linewidth=1.0,
                                alpha=0.75, zorder=5))
            ax.add_patch(Circle((lon, lat), 0.9, color="#F0E0C0", zorder=6))
        elif role == "dream":
            for r, a in [(4.5, 0.10), (3.0, 0.18)]:
                ax.add_patch(Circle((lon, lat), r, color="#E8D4A8",
                                    alpha=a, zorder=5))
            ax.add_patch(Circle((lon, lat), 3.2, fill=False,
                                edgecolor="#E8D4A8", linewidth=0.8,
                                alpha=0.55, zorder=5))
            ax.add_patch(Circle((lon, lat), 1.8, fill=False,
                                edgecolor="#E8D4A8", linewidth=0.9,
                                alpha=0.85, zorder=5))
            ax.add_patch(Circle((lon, lat), 0.8, color="#F5E6C0", zorder=6))
        else:  # present
            for r, a in [(2.0, 0.12), (1.3, 0.22)]:
                ax.add_patch(Circle((lon, lat), r, color="#D4B574",
                                    alpha=a, zorder=5))
            ax.add_patch(Circle((lon, lat), 0.8, color="#D4B574", zorder=6))

        # leader line do ponto ate o label
        label_x, label_y, ha = LABEL_POS[label]
        ax.plot([lon, label_x], [lat, label_y],
                color="#8A7342", linewidth=0.5, alpha=0.7, zorder=4)

        # cor/tamanho do label
        if role == "origin":
            color, fs, fw = "#F5E6C0", 11, "bold"
        elif role == "next":
            color, fs, fw = "#F0E0C0", 10, "bold"
        elif role == "dream":
            color, fs, fw = "#F5E6C0", 10, "bold"
        else:
            color, fs, fw = "#D4B574", 9, "normal"

        # nudge do texto pra nao colar no fim do leader
        if ha == "right":
            tx = label_x - 1.5
        else:
            tx = label_x + 1.5
        ax.text(
            tx, label_y, label,
            color=color, fontsize=fs, fontfamily="Helvetica",
            fontweight=fw, ha=ha, va="center", zorder=7,
        )

    # (G) Vinheta nos cantos — 4 elipses pretas pra dar foco no centro
    for cx, cy in [(-200, 90), (200, 90), (-200, -75), (200, -75)]:
        ax.add_patch(Ellipse((cx, cy), width=200, height=120,
                              facecolor="#0A0A0A", edgecolor="none",
                              alpha=0.7, zorder=10))

    plt.tight_layout(pad=0)
    plt.savefig(out, facecolor="#0A0A0A", dpi=220, bbox_inches="tight", pad_inches=0.1)
    plt.close(fig)
    print(f"  mapa ok: {out}")
    return out

WORLD_MAP = build_world_map()

# -----------------------------------------------------------------------------
# Gerar overlays de gradiente (preto solido -> transparente) como PNG
# -----------------------------------------------------------------------------
def make_gradient_png(name, w, h, direction="vertical", start_alpha=255, end_alpha=0):
    """Cria PNG com gradiente preto transparente. Salva em IMG_DIR."""
    out = f"{IMG_DIR}/grad-{name}.png"
    img = Image.new("RGBA", (w, h), (0, 0, 0, 0))
    px = img.load()
    if direction == "vertical":
        for y in range(h):
            t = y / max(h - 1, 1)
            a = int(start_alpha + (end_alpha - start_alpha) * t)
            for x in range(w):
                px[x, y] = (0, 0, 0, a)
    else:  # horizontal
        for x in range(w):
            t = x / max(w - 1, 1)
            a = int(start_alpha + (end_alpha - start_alpha) * t)
            for y in range(h):
                px[x, y] = (0, 0, 0, a)
    img.save(out, "PNG", optimize=True)
    return out

# Pre-rendering: para cada slide com fullbleed, vamos sobrepor um overlay solido
# e um gradiente. Gerar uma vez so:
GRAD_BOTTOM = make_gradient_png("bottom", 1600, 900, "vertical", 255, 0)  # preto embaixo -> transp em cima
# para usar invertido, podemos rotacionar/flipar — mas mais simples gerar outro
GRAD_TOP = make_gradient_png("top", 1600, 900, "vertical", 0, 255)

# -----------------------------------------------------------------------------
# Helpers de slide
# -----------------------------------------------------------------------------
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = prs.slide_width
SH = prs.slide_height
BLANK = prs.slide_layouts[6]

def add_bg(slide, color=BLACK):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.shadow.inherit = False
    # remover shadow
    sp = bg.shadow
    return bg

def add_image_fullbleed(slide, path):
    """Coloca imagem cobrindo todo o slide, recortando se necessario.
    python-pptx nao tem crop facil, entao usamos a imagem inteira escalada para
    cobrir e centralizada. Para webp ja convertido, garantimos cobertura."""
    # Descobrir aspecto e ajustar para cover
    with Image.open(path) as im:
        iw, ih = im.size
    slide_ratio = SW / SH
    img_ratio = iw / ih
    if img_ratio > slide_ratio:
        # imagem mais larga: ajusta altura, sobra lateral
        new_h = SH
        new_w = int(new_h * img_ratio)
        left = (SW - new_w) // 2
        top = 0
    else:
        new_w = SW
        new_h = int(new_w / img_ratio)
        left = 0
        top = (SH - new_h) // 2
    slide.shapes.add_picture(path, left, top, width=new_w, height=new_h)

def add_overlay(slide, alpha_pct):
    """Retangulo preto solido com transparencia. alpha_pct = 0..100 (cobertura)."""
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    rect.line.fill.background()
    rect.fill.solid()
    rect.fill.fore_color.rgb = BLACK
    # Aplicar transparencia via XML (python-pptx nao expoe diretamente)
    sp = rect.fill._xPr
    # Inserir alpha no solidFill
    solid_fill = rect.fill._xPr.find(qn("a:solidFill"))
    if solid_fill is None:
        # fallback: encontrar via xpath
        from pptx.oxml.ns import nsmap
        solid_fill = rect.fill._xPr.find(".//" + qn("a:solidFill"))
    srgb = solid_fill.find(qn("a:srgbClr"))
    # adiciona alpha child
    alpha_val = int(alpha_pct * 1000)  # 0..100000 (60% -> 60000 opacity)
    # No OOXML, a:alpha val=100000 = 100% opaque. Queremos cobertura = alpha_pct
    alpha_el = etree.SubElement(srgb, qn("a:alpha"))
    alpha_el.set("val", str(alpha_val))
    return rect

def add_gradient_overlay(slide, path, height_frac=1.0, anchor="bottom"):
    """Insere gradiente como imagem. height_frac controla altura relativa."""
    h = int(SH * height_frac)
    if anchor == "bottom":
        top = SH - h
    else:
        top = 0
    slide.shapes.add_picture(path, 0, top, width=SW, height=h)

def add_text(slide, left, top, width, height, text, *,
             font=F_SERIF, size=24, color=CREAM, bold=False, italic=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, letter_spacing=None,
             line_spacing=1.15):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    # quebrar em linhas usando <br>
    lines = text.split("<br>")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        f = r.font
        f.name = font
        f.size = Pt(size)
        f.bold = bold
        f.italic = italic
        f.color.rgb = color
        if letter_spacing is not None:
            # OOXML: rPr spc em centesimos de ponto. positivo amplia.
            rPr = r._r.get_or_add_rPr()
            rPr.set("spc", str(int(letter_spacing * 100)))
    return tb

def add_hline(slide, left, top, width, color=GOLD, weight=0.75):
    line = slide.shapes.add_connector(1, left, top, left + width, top)
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    return line

def add_vline(slide, left, top, height, color=GOLD, weight=0.5):
    line = slide.shapes.add_connector(1, left, top, left, top + height)
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    return line

def add_dot(slide, cx, cy, r_emu, color=GOLD):
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                 cx - r_emu, cy - r_emu, r_emu * 2, r_emu * 2)
    dot.line.fill.background()
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    return dot

# -----------------------------------------------------------------------------
# SLIDE 1 — CAPA
# -----------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_bg(s, BLACK)
add_image_fullbleed(s, HERO_JPG)
add_overlay(s, 55)  # 55% preto
# gradiente forte embaixo (preto solido) -> transparente em cima
add_gradient_overlay(s, GRAD_BOTTOM, height_frac=0.7, anchor="bottom")

# Wordmark dourado pequeno no topo central
WM_W = Inches(2.0)
WM_H = Inches(0.5)
s.shapes.add_picture(WORDMARK_GOLD,
                     (SW - WM_W) // 2, Inches(0.45),
                     width=WM_W, height=WM_H)

# Eyebrow
add_text(s, Inches(0), Inches(2.4), SW, Inches(0.4),
         "VALIM FARMS  ·  SÃO PAULO  ·  BRAZIL",
         font=F_SANS, size=11, color=GOLD_LIGHT, align=PP_ALIGN.CENTER,
         letter_spacing=5, bold=False)

# Headline removido — apresentador escreve a frase ao vivo / na hora
# (mantemos somente wordmark, eyebrow, linha dourada e atribuicao)

# Linha dourada fina
add_hline(s, Inches(5.5), Inches(6.35), Inches(2.333), GOLD, 0.75)

# Atribuicao embaixo
add_text(s, Inches(0), Inches(6.55), SW, Inches(0.4),
         "WILSON LUIZ VALIM ZERBINATTI  ·  THIRD GENERATION  ·  1897",
         font=F_SANS, size=10, color=GOLD, align=PP_ALIGN.CENTER,
         letter_spacing=4)

# -----------------------------------------------------------------------------
# SLIDE 2 — LINHAGEM
# -----------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_bg(s, BLACK)

# Esquerda: imagem ~50%
img_w = Inches(6.0)
img_h = SH
# Inserir imagem flowers-sunset cobrindo metade esquerda
# fazer cover manual
with Image.open(JPGS["flowers-sunset"]) as im:
    iw, ih = im.size
target_ratio = img_w / img_h
img_ratio = iw / ih
if img_ratio > target_ratio:
    new_h = img_h
    new_w = int(new_h * img_ratio)
    px_left = -(new_w - img_w) // 2
    px_top = 0
else:
    new_w = img_w
    new_h = int(new_w / img_ratio)
    px_left = 0
    px_top = (img_h - new_h) // 2
# precisamos recortar a area visivel manualmente — mais simples: gerar JPG cropado
def crop_cover(src_path, target_w_emu, target_h_emu, out_name):
    out = f"{IMG_DIR}/crop-{out_name}.jpg"
    target_ratio = target_w_emu / target_h_emu
    with Image.open(src_path) as im:
        iw, ih = im.size
        ir = iw / ih
        if ir > target_ratio:
            new_w = int(ih * target_ratio)
            left = (iw - new_w) // 2
            box = (left, 0, left + new_w, ih)
        else:
            new_h = int(iw / target_ratio)
            top = (ih - new_h) // 2
            box = (0, top, iw, top + new_h)
        im2 = im.crop(box)
        im2.save(out, "JPEG", quality=92, optimize=True)
    return out

crop1 = crop_cover(JPGS["flowers-sunset"], img_w, img_h, "s2-left")
s.shapes.add_picture(crop1, 0, 0, width=img_w, height=img_h)

# overlay sutil escuro na imagem para integrar com preto
ov = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, img_w, SH)
ov.line.fill.background()
ov.fill.solid()
ov.fill.fore_color.rgb = BLACK
srgb = ov.fill._xPr.find(".//" + qn("a:srgbClr"))
alpha_el = etree.SubElement(srgb, qn("a:alpha"))
alpha_el.set("val", "25000")  # 25% cobertura

# Lado direito: texto editorial
RX = Inches(6.7)
RW = Inches(6.0)

# Eyebrow
add_text(s, RX, Inches(1.2), RW, Inches(0.4),
         "— LINEAGE",
         font=F_SANS, size=11, color=GOLD, letter_spacing=5)

# Headline
add_text(s, RX, Inches(1.7), RW, Inches(2.7),
         "Three generations.<br>One land.<br>One coffee.",
         font=F_SERIF, size=54, color=CREAM, italic=True, line_spacing=1.1)

# Divisor curto
add_hline(s, RX, Inches(4.55), Inches(0.7), GOLD, 0.75)

# Corpo
body = ("An Italian family arrived in São Paulo in 1897. They planted. They waited. "
        "They learned. Today I am Wilson Luiz Valim Zerbinatti — the third to walk these "
        "coffee fields. The Zerbinatti surname became a brand, but it\u2019s still what it "
        "always was: people who listen to the land.")
add_text(s, RX, Inches(4.85), RW, Inches(1.6),
         body, font=F_SANS, size=13, color=CREAM, line_spacing=1.4)

# Stat destaque embaixo
add_hline(s, RX, Inches(6.45), Inches(5.5), GOLD, 0.5)
add_text(s, RX, Inches(6.6), RW, Inches(0.5),
         "129 YEARS  ·  THREE GENERATIONS  ·  ONE ORIGIN",
         font=F_SANS, size=11, color=GOLD_LIGHT, letter_spacing=4, bold=True)

# -----------------------------------------------------------------------------
# SLIDE 3 — A TERRA (tabela editorial + cerejas)
# -----------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_bg(s, BLACK)

# Imagem direita ~60%
right_w = Inches(8.0)
right_h = SH
crop2 = crop_cover(JPGS["cherries-yellow"], right_w, right_h, "s3-right")
s.shapes.add_picture(crop2, SW - right_w, 0, width=right_w, height=right_h)

# Gradiente para esquerda (fade da imagem pra preto)
grad_lr = make_gradient_png("lr-s3", 1200, 900, "horizontal", 255, 0)
s.shapes.add_picture(grad_lr, SW - right_w - Inches(1.5), 0,
                     width=Inches(3.5), height=SH)

# overlay sutil na imagem
ov = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, SW - right_w, 0, right_w, SH)
ov.line.fill.background()
ov.fill.solid()
ov.fill.fore_color.rgb = BLACK
srgb = ov.fill._xPr.find(".//" + qn("a:srgbClr"))
alpha_el = etree.SubElement(srgb, qn("a:alpha"))
alpha_el.set("val", "30000")

LX = Inches(0.7)
LW = Inches(5.5)

add_text(s, LX, Inches(1.0), LW, Inches(0.4),
         "— THE LAND",
         font=F_SANS, size=11, color=GOLD, letter_spacing=5)

add_text(s, LX, Inches(1.5), LW, Inches(2.2),
         "Where everything<br>happens on its own.",
         font=F_SERIF, size=52, color=CREAM, italic=True, line_spacing=1.05)

# Tabela editorial
table_data = [
    ("LOCATION", "Valim Farms, SP"),
    ("ALTITUDE", "640 – 760 m"),
    ("SOIL", "Volcanic"),
    ("VARIETIES", "Arara · Paraiso 2"),
    ("PROCESS", "Natural, sun-dried"),
    ("HARVEST", "Hand-picked, ripe cherry"),
]
T_TOP = Inches(4.0)
ROW_H = Inches(0.45)
COL1_W = Inches(2.0)
COL2_W = Inches(3.5)
for i, (k, v) in enumerate(table_data):
    y = T_TOP + ROW_H * i
    # linha de cima
    add_hline(s, LX, y, COL1_W + COL2_W, GOLD_DARK, 0.5)
    # coluna 1 — label mono caps ouro
    add_text(s, LX, y + Inches(0.10), COL1_W, ROW_H,
             k, font=F_SANS, size=9.5, color=GOLD,
             letter_spacing=4, bold=True)
    # coluna 2 — valor cream
    add_text(s, LX + COL1_W, y + Inches(0.08), COL2_W, ROW_H,
             v, font=F_SERIF, size=15, color=CREAM, italic=True)
# linha final
add_hline(s, LX, T_TOP + ROW_H * len(table_data), COL1_W + COL2_W,
          GOLD_DARK, 0.5)

# Tira de selos no rodape do lado esquerdo (credibilidade institucional)
SELOS = [
    f"{ROOT}/public/assets/selo-cup.png",
    f"{ROOT}/public/assets/selo-sca-90-50.png",
    f"{ROOT}/public/assets/selo-organico-fazenda.png",
    f"{ROOT}/public/assets/selo-brasil.png",
]
SELO_SIZE = Inches(0.55)
SELO_AREA_LEFT = Inches(0.7)
SELO_AREA_RIGHT = Inches(6.2)
SELO_AREA_W = SELO_AREA_RIGHT - SELO_AREA_LEFT
SELO_Y = Inches(6.85)
# eyebrow acima dos selos
add_text(s, SELO_AREA_LEFT, Inches(6.55), SELO_AREA_W, Inches(0.3),
         "— CERTIFICATIONS",
         font=F_SANS, size=9, color=GOLD_DARK, letter_spacing=4, bold=True)
# espacamento uniforme entre centros dos selos
n_selos = len(SELOS)
slot_w = SELO_AREA_W // n_selos
for idx, selo_path in enumerate(SELOS):
    center_x = SELO_AREA_LEFT + slot_w * idx + slot_w // 2
    selo_left = center_x - SELO_SIZE // 2
    s.shapes.add_picture(selo_path, selo_left, SELO_Y,
                         width=SELO_SIZE, height=SELO_SIZE)

# -----------------------------------------------------------------------------
# SLIDE 4 — BRAND PORTFOLIO / VALIM FARMS · TERROIR (EN)
# -----------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_bg(s, BLACK)

# TOPO: eyebrow esquerdo + italic ouro direito
add_text(s, Inches(0.7), Inches(0.55), Inches(8), Inches(0.4),
         "ZERBINATTI  ·  BRAND PORTFOLIO",
         font=F_SANS, size=9, color=GOLD_DARK, letter_spacing=4, bold=True)

add_text(s, Inches(9.5), Inches(0.55), Inches(3.2), Inches(0.4),
         "Famiglia · since 1897",
         font=F_SERIF, size=12, color=GOLD, italic=True,
         align=PP_ALIGN.RIGHT)

# HEADLINE + BODY
add_text(s, Inches(0.7), Inches(1.3), Inches(8), Inches(0.4),
         "— VALIM FARMS  ·  TERROIR",
         font=F_SANS, size=10, color=GOLD, letter_spacing=5, bold=True)

add_text(s, Inches(0.7), Inches(1.7), Inches(6.5), Inches(1.5),
         "One farm.<br>A whole climate.",
         font=F_SERIF, size=44, color=CREAM, italic=True, line_spacing=1.05)

s4_body_en = ("Valim Farms sits on natural terraces between 640 and 760 meters, "
              "with deep volcanic soil, dry winters and a long, slow ripening window. "
              "Two Arabica varieties — Arara and Paraíso 2 — are planted by hand and "
              "harvested only when the cherry is fully ripe.")
add_text(s, Inches(7.5), Inches(1.85), Inches(5.1), Inches(1.3),
         s4_body_en, font=F_SANS, size=11, color=CREAM, line_spacing=1.45)

# LINHA DIVISORA
add_hline(s, Inches(0.7), Inches(3.4), Inches(11.93), GOLD_DARK, 0.5)

# GRID FOTOS
# Coluna A — foto principal grande
s4_w_a = Inches(4.2)
s4_h_a = Inches(3.2)
s4_crop_a = crop_cover(JPGS["cherries-yellow"], s4_w_a, s4_h_a, "s4-photo-a")
s.shapes.add_picture(s4_crop_a, Inches(0.7), Inches(3.7),
                     width=s4_w_a, height=s4_h_a)

# Coluna B — duas fotos empilhadas
s4_w_b = Inches(3.5)
s4_h_b = Inches(1.55)
s4_crop_b1 = crop_cover(JPGS["young-plants"], s4_w_b, s4_h_b, "s4-photo-b1")
s.shapes.add_picture(s4_crop_b1, Inches(5.0), Inches(3.7),
                     width=s4_w_b, height=s4_h_b)
s4_crop_b2 = crop_cover(JPGS["drying-yard-tree"], s4_w_b, s4_h_b, "s4-photo-b2")
s.shapes.add_picture(s4_crop_b2, Inches(5.0), Inches(5.35),
                     width=s4_w_b, height=s4_h_b)

# Coluna C — CARD TERROIR
card = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                          Inches(8.7), Inches(3.7),
                          Inches(3.93), Inches(3.2))
card.fill.solid()
card.fill.fore_color.rgb = BLACK
card.line.color.rgb = GOLD_DARK
card.line.width = Pt(0.5)

# Headline do card
add_text(s, Inches(8.95), Inches(3.95), Inches(3.4), Inches(0.5),
         "The terroir.",
         font=F_SERIF, size=22, color=CREAM, italic=True)

# Body do card
s4_card_body_en = ("Volcanic soil, terraced slopes, cool mornings, dry winters. "
                   "The land does most of the work — we just stay out of its way.")
add_text(s, Inches(8.95), Inches(4.55), Inches(3.4), Inches(1.4),
         s4_card_body_en, font=F_SANS, size=10, color=CREAM, line_spacing=1.45)

# Linha divisora dentro do card
add_hline(s, Inches(8.95), Inches(5.85), Inches(3.4), GOLD_DARK, 0.5)

# Stats grid 2x2
# Esquerda top: 760m / PEAK ALTITUDE
add_text(s, Inches(8.95), Inches(6.05), Inches(1.7), Inches(0.4),
         "760m", font=F_SERIF, size=20, color=GOLD, italic=True)
add_text(s, Inches(8.95), Inches(6.40), Inches(1.7), Inches(0.3),
         "PEAK ALTITUDE", font=F_SANS, size=8, color=GOLD_DARK,
         letter_spacing=4, bold=True)

# Direita top: 02 / VARIETIES
add_text(s, Inches(10.75), Inches(6.05), Inches(1.7), Inches(0.4),
         "02", font=F_SERIF, size=20, color=GOLD, italic=True)
add_text(s, Inches(10.75), Inches(6.40), Inches(1.7), Inches(0.3),
         "VARIETIES", font=F_SANS, size=8, color=GOLD_DARK,
         letter_spacing=4, bold=True)

# Esquerda bottom: 100% / ARABICA
add_text(s, Inches(8.95), Inches(6.55), Inches(1.7), Inches(0.4),
         "100%", font=F_SERIF, size=20, color=GOLD, italic=True)
add_text(s, Inches(8.95), Inches(6.90), Inches(1.7), Inches(0.3),
         "ARABICA", font=F_SANS, size=8, color=GOLD_DARK,
         letter_spacing=4, bold=True)

# Direita bottom: May-Sep / HARVEST
add_text(s, Inches(10.75), Inches(6.55), Inches(1.7), Inches(0.4),
         "May-Sep", font=F_SERIF, size=20, color=GOLD, italic=True)
add_text(s, Inches(10.75), Inches(6.90), Inches(1.7), Inches(0.3),
         "HARVEST", font=F_SANS, size=8, color=GOLD_DARK,
         letter_spacing=4, bold=True)

# -----------------------------------------------------------------------------
# SLIDE 5 — O MUNDO (mapa)
# -----------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_bg(s, BLACK)

# Eyebrow
add_text(s, Inches(0.7), Inches(0.7), Inches(12), Inches(0.4),
         "— THE COFFEE ALREADY CROSSES OCEANS",
         font=F_SANS, size=11, color=GOLD, letter_spacing=5, bold=True)

# Headline
add_text(s, Inches(0.7), Inches(1.15), Inches(12), Inches(1.6),
         "Three flags today.<br>The dream is Oregon.",
         font=F_SERIF, size=42, color=CREAM, italic=True, line_spacing=1.1)

# Mapa central
map_w = Inches(11.5)
map_h = Inches(4.2)
map_left = (SW - map_w) // 2
map_top = Inches(2.8)
s.shapes.add_picture(WORLD_MAP, map_left, map_top, width=map_w, height=map_h)

# Microcopy embaixo
add_text(s, Inches(0), Inches(7.0), SW, Inches(0.4),
         "PRESENCE  ·  EXPANSION  ·  2026",
         font=F_SANS, size=11, color=GOLD_LIGHT,
         letter_spacing=5, align=PP_ALIGN.CENTER, bold=True)

# -----------------------------------------------------------------------------
# SLIDE 6 — OS NUMEROS
# -----------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_bg(s, BLACK)

# Sutil radial: criar via gradient PNG redondo
def make_radial(name, w=1600, h=900, color=(201, 169, 97), max_alpha=70):
    out = f"{IMG_DIR}/radial-{name}.png"
    img = Image.new("RGBA", (w, h), (0, 0, 0, 0))
    px = img.load()
    cx, cy = w / 2, h / 2
    rmax = ((w / 2) ** 2 + (h / 2) ** 2) ** 0.5 * 0.7
    for y in range(h):
        for x in range(w):
            d = ((x - cx) ** 2 + (y - cy) ** 2) ** 0.5
            t = min(d / rmax, 1.0)
            a = int(max_alpha * (1 - t) ** 2)
            px[x, y] = (color[0], color[1], color[2], a)
    img.save(out, "PNG", optimize=True)
    return out

radial = make_radial("s5", 1200, 700, (201, 169, 97), 55)
s.shapes.add_picture(radial, 0, 0, width=SW, height=SH)

# Eyebrow
add_text(s, Inches(0.7), Inches(0.8), Inches(12), Inches(0.4),
         "— THE PROOF",
         font=F_SANS, size=11, color=GOLD, letter_spacing=5, bold=True)

# Headline
add_text(s, Inches(0.7), Inches(1.3), Inches(12), Inches(1.6),
         "The numbers<br>that matter.",
         font=F_SERIF, size=46, color=CREAM, italic=True, line_spacing=1.1)

# HERO NUMBER (esquerda) — placeholder gigante de sacas/ano
HERO_LEFT = Inches(0.7)
HERO_TOP = Inches(3.5)
HERO_W = Inches(6.5)
add_text(s, HERO_LEFT, HERO_TOP, HERO_W, Inches(2.4),
         "[ ___ ]", font=F_SERIF, size=160, color=GOLD_DARK,
         italic=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
         line_spacing=1.0)
add_text(s, HERO_LEFT, HERO_TOP + Inches(2.5), HERO_W, Inches(0.4),
         "BAGS / YEAR", font=F_SANS, size=12, color=CREAM,
         letter_spacing=5, align=PP_ALIGN.CENTER, bold=True)
add_text(s, HERO_LEFT, HERO_TOP + Inches(2.95), HERO_W, Inches(0.4),
         "2026 production", font=F_SERIF, size=13, color=GOLD_LIGHT,
         italic=True, align=PP_ALIGN.CENTER)

# linha vertical divisoria
add_vline(s, Inches(7.2), Inches(3.7), Inches(3.2), GOLD_DARK, 0.5)

# 3 STATS MENORES (direita), empilhados
STATS_X = Inches(7.5)
STATS_W = Inches(5.3)
STATS = [
    ("20 CONTAINERS",  "SPAIN SHIPMENT · 2026",     "value"),
    ("[ __ ]",         "SCAA SCORE",                "placeholder"),
    ("100%",           "TRACEABLE MICROLOT",        "value"),
]
SUBLABELS = [
    "Scheduled shipment",
    "Cup rating",
    "Origin guaranteed",
]
STATS_TOP = Inches(3.7)
STATS_BLOCK_H = Inches(1.07)
for i, ((num, label, kind), sub) in enumerate(zip(STATS, SUBLABELS)):
    y = STATS_TOP + STATS_BLOCK_H * i
    color_num = GOLD_DARK if kind == "placeholder" else GOLD
    add_text(s, STATS_X, y, STATS_W, Inches(0.6),
             num, font=F_SERIF, size=38, color=color_num,
             italic=True, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             line_spacing=1.0)
    add_text(s, STATS_X, y + Inches(0.6), STATS_W, Inches(0.3),
             label, font=F_SANS, size=10, color=GOLD,
             letter_spacing=4, bold=True, align=PP_ALIGN.LEFT)
    add_text(s, STATS_X, y + Inches(0.88), STATS_W, Inches(0.3),
             sub, font=F_SERIF, size=11, color=GOLD_LIGHT,
             italic=True, align=PP_ALIGN.LEFT)

# -----------------------------------------------------------------------------
# SLIDE 7 — ESTADOS UNIDOS
# -----------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_bg(s, BLACK)

add_image_fullbleed(s, JPGS["drying-yard-tree"])
add_overlay(s, 72)  # 72% preto

# Gradiente vertical forte
add_gradient_overlay(s, GRAD_BOTTOM, height_frac=1.0, anchor="bottom")
# overlay extra escuro
ov = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
ov.line.fill.background()
ov.fill.solid()
ov.fill.fore_color.rgb = BLACK
srgb = ov.fill._xPr.find(".//" + qn("a:srgbClr"))
alpha_el = etree.SubElement(srgb, qn("a:alpha"))
alpha_el.set("val", "30000")

# Eyebrow topo
add_text(s, Inches(0), Inches(0.7), SW, Inches(0.4),
         "— MY NEXT DESTINATION",
         font=F_SANS, size=11, color=GOLD, letter_spacing=5,
         align=PP_ALIGN.CENTER, bold=True)

# Headline GIGANTE
add_text(s, Inches(0.3), Inches(1.7), Inches(12.7), Inches(3.3),
         "United States.",
         font=F_SERIF, size=180, color=GOLD, italic=True,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)

# linha dourada decorativa
add_hline(s, Inches(5.0), Inches(5.5), Inches(3.333), GOLD, 1.0)

# Subheadline
add_text(s, Inches(1.5), Inches(5.7), Inches(10.333), Inches(1.2),
         "I would stop exporting to Europe<br>to finally arrive here.",
         font=F_SERIF, size=28, color=CREAM, italic=True,
         align=PP_ALIGN.CENTER, line_spacing=1.25)

# Atribuicao
add_text(s, Inches(0), Inches(7.05), SW, Inches(0.35),
         "WILSON LUIZ VALIM ZERBINATTI  ·  2026",
         font=F_SANS, size=10, color=GOLD_LIGHT, letter_spacing=4,
         align=PP_ALIGN.CENTER, bold=True)

# -----------------------------------------------------------------------------
# SLIDE 8 — ENCERRAMENTO
# -----------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_bg(s, BLACK)
add_image_fullbleed(s, JPGS["peneirar"])
add_overlay(s, 78)

# Frase central com aspas tipograficas
add_text(s, Inches(1.0), Inches(0.9), Inches(11.333), Inches(2.3),
         "\u201CThe land teaches, the sun decides,<br>"
         "we just observe.\u201D",
         font=F_SERIF, size=36, color=CREAM, italic=True,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.3)

# Ornamento dourado: linha + diamante centro
orn_y = Inches(3.45)
add_hline(s, Inches(4.0), orn_y, Inches(2.3), GOLD, 0.75)
add_hline(s, Inches(7.0), orn_y, Inches(2.3), GOLD, 0.75)
# diamante central
diamond_r = Inches(0.08)
diamond = s.shapes.add_shape(MSO_SHAPE.DIAMOND,
                              Inches(6.667) - diamond_r,
                              orn_y - diamond_r,
                              diamond_r * 2, diamond_r * 2)
diamond.line.color.rgb = GOLD
diamond.line.width = Pt(0.5)
diamond.fill.solid()
diamond.fill.fore_color.rgb = GOLD

# Atribuicao
add_text(s, Inches(0), Inches(3.7), SW, Inches(0.4),
         "— WILSON LUIZ VALIM ZERBINATTI  ·  THIRD GENERATION",
         font=F_SANS, size=9, color=GOLD, letter_spacing=5,
         align=PP_ALIGN.CENTER, bold=True)

# QR code centralizado (medalhao redondo)
QR_SIZE = Inches(2.1)
qr_x = (SW - QR_SIZE) / 2
qr_y = Inches(4.4)
s.shapes.add_picture(QR_PATH, qr_x, qr_y, width=QR_SIZE, height=QR_SIZE)

# Rodape: Instagram + site abaixo do QR
add_text(s, Inches(0), Inches(6.75), SW, Inches(0.4),
         "@ZERBINATTICOFFEE  ·  ZERBINATTICOFFEE.COM",
         font=F_SANS, size=12, color=GOLD_LIGHT, letter_spacing=3,
         align=PP_ALIGN.CENTER, bold=True)

# Wordmark dourado canto inferior direito
WM2_W = Inches(1.2)
WM2_H = Inches(0.3)
s.shapes.add_picture(WORDMARK_GOLD,
                     SW - WM2_W - Inches(0.4),
                     SH - WM2_H - Inches(0.25),
                     width=WM2_W, height=WM2_H)

# -----------------------------------------------------------------------------
# SALVAR
# -----------------------------------------------------------------------------
prs.save(OUT_PATH)
size_kb = os.path.getsize(OUT_PATH) / 1024
print(f"\nOK -> {OUT_PATH}")
print(f"slides: {len(prs.slides)}  size: {size_kb:.1f} KB")
