from pptx import Presentation
from pptx.util import Emu
import os

PATH = "/Users/fabricio/dev/cafe/docs/palestra/zerbinatti-palestra-5min.pptx"
prs = Presentation(PATH)
print(f"Slides: {len(prs.slides)}  size: {prs.slide_width/914400:.2f}x{prs.slide_height/914400:.2f} in")
print(f"File: {os.path.getsize(PATH)/1024:.1f} KB\n")

for i, s in enumerate(prs.slides, 1):
    print(f"--- SLIDE {i} ({len(s.shapes)} shapes) ---")
    for sh in s.shapes:
        kind = sh.shape_type
        name = sh.name
        text = ""
        if sh.has_text_frame:
            txt = sh.text_frame.text.replace("\n", " | ")
            if len(txt) > 90:
                txt = txt[:87] + "..."
            text = f"  TEXT: {txt!r}"
        print(f"  [{kind}] {name}{text}")
    print()
